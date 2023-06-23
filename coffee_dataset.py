import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import os
import openai
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

# Set a random seed for reproducibility
np.random.seed(0)

# Create a DataFrame
data = pd.DataFrame({
    'Location ID': range(1, 6),
    'Total Sales': np.random.randint(1e5, 5e5, 5),
    'Average Customer Rating': np.random.uniform(1, 5, 5).round(1),
    'Number of Employees': np.random.randint(5, 50, 5),
    'Operating Costs': np.random.randint(1e4, 5e4, 5)
})

# Save the DataFrame to a CSV file
data.to_csv('coffee_chain_data.csv', index=False)

# Load the data from CSV file
data = pd.read_csv('coffee_chain_data.csv')

# Create a directory for the output
os.makedirs('Final_Report', exist_ok=True)

# Scatter plot of Total Sales vs Average Customer Rating
plt.scatter(data['Total Sales'], data['Average Customer Rating'])
plt.xlabel('Total Sales')
plt.ylabel('Average Customer Rating')
plt.title('Scatter plot of Total Sales vs Average Customer Rating')
plt.savefig('output/plot1.png')  # Save the plot as an image
plt.show()

# Scatter plot of Number of Employees vs Operating Costs
plt.scatter(data['Number of Employees'], data['Operating Costs'])
plt.xlabel('Number of Employees')
plt.ylabel('Operating Costs')
plt.title('Scatter plot of Number of Employees vs Operating Costs')
plt.savefig('output/plot2.png')  # Save the plot as an image
plt.show()

openai.api_key = os.getenv('OPENAI_KEY')

# Create a chat models prompt
message1 = openai.ChatCompletion.create(
  model="gpt-3.5-turbo-0613",
  messages=[
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": f"The profits for the five stores in our coffee chain are as follows: {data['Total Sales'].tolist()}. Could you please provide a detailed analysis and comparison of the performance of these stores?"},
    ]
)

# Print the assistant's response
print(message1['choices'][0]['message']['content'])

# Now create a PDF report
pdf = FPDF()

# Add a page
pdf.add_page()

# Set font
pdf.set_font("Arial", size = 15)

# Add a cell with the text from GPT-3
pdf.cell(200, 10, txt=message1['choices'][0]['message']['content'], ln=True, align='C')

# Add images of the plots
pdf.image('output/plot1.png', x=10, y=30, w=100)
pdf.image('output/plot2.png', x=10, y=140, w=100)

# Save the pdf with name .pdf
pdf.output("output/report.pdf")

# Creating a presentation
prs = Presentation()

# Adding a slide layout
slide_layout = prs.slide_layouts[5]  # This is a blank slide

# Add a slide with text
slide = prs.slides.add_slide(slide_layout)
txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
tf = txBox.text_frame
tf.text = message1['choices'][0]['message']['content']

# Add slides with images
for plot in ['plot1.png', 'plot2.png']:
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(f'output/{plot}', Inches(1), Inches(1), Inches(6), Inches(4.5))

# Save the presentation
prs.save('output/presentation.pptx')
